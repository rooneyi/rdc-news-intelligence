from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x0D, 0x2B, 0x55)   # #0D2B55  fond titre
BLEU_CLAIR  = RGBColor(0x1A, 0x73, 0xE8)   # #1A73E8  accents
VERT        = RGBColor(0x0F, 0x96, 0x6B)   # #0F966B  bullet points
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR  = RGBColor(0xF4, 0xF6, 0xFB)
GRIS_TEXTE  = RGBColor(0x33, 0x33, 0x33)
ORANGE      = RGBColor(0xE8, 0x71, 0x1A)   # mise en évidence

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout vide


# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, italic=False,
             color=BLANC, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   size=Pt(15), color=GRIS_TEXTE,
                   bullet_color=VERT, bullet="▶  "):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        # bullet run
        br = p.add_run()
        br.text = bullet
        br.font.size = size
        br.font.color.rgb = bullet_color
        br.font.bold = True
        # text run
        tr = p.add_run()
        tr.text = item
        tr.font.size = size
        tr.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.4, fill_rgb=BLEU_FONCE)
    add_text(slide, title, 0.3, 0.1, 12, 0.75,
             size=Pt(28), bold=True, color=BLANC)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.85, 12, 0.45,
                 size=Pt(14), color=RGBColor(0xAA, 0xCC, 0xFF))
    # thin accent line
    add_rect(slide, 0, 1.4, 13.33, 0.06, fill_rgb=BLEU_CLAIR)


def footer(slide, num, total=8):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=BLEU_FONCE)
    add_text(slide, "RDC News Intelligence  |  KALUMBA ILUNGA Rooney  |  UDBL – Génie Logiciel 2026",
             0.3, 7.17, 10, 0.28, size=Pt(9), color=RGBColor(0xAA, 0xCC, 0xFF))
    add_text(slide, f"{num} / {total}", 12.5, 7.17, 0.7, 0.28,
             size=Pt(9), color=BLANC, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – PAGE DE TITRE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=BLEU_FONCE)
# Accent dégradé simulé via bandes
for i in range(8):
    add_rect(sl, 0, 7.5 - (i+1)*0.35, 13.33, 0.35,
             fill_rgb=RGBColor(0x0D+i*3, 0x2B+i*5, 0x55+i*8))

# Logo carré coloré (placeholder)
add_rect(sl, 0.4, 0.4, 1.2, 1.2, fill_rgb=BLEU_CLAIR)
add_text(sl, "UDBL", 0.42, 0.55, 1.15, 0.6,
         size=Pt(20), bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# Université
add_text(sl, "UNIVERSITÉ DON BOSCO DE LUBUMBASHI",
         1.8, 0.3, 11, 0.45, size=Pt(13), color=RGBColor(0xAA, 0xCC, 0xFF))
add_text(sl, "Faculté des Sciences Informatiques  ·  Génie Logiciel",
         1.8, 0.72, 11, 0.4, size=Pt(11), italic=True,
         color=RGBColor(0x88, 0xAA, 0xDD))

# Ligne séparatrice
add_rect(sl, 0.4, 1.35, 12.5, 0.06, fill_rgb=BLEU_CLAIR)

# Titre principal
add_text(sl, "DÉVELOPPEMENT D'UN SYSTÈME INTELLIGENT DE\n"
             "RECOMMANDATION ET DE RECHERCHE D'INFORMATION",
         0.4, 1.6, 12.5, 1.2,
         size=Pt(26), bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text(sl, "pour lutter contre la surinformation et la désinformation en ligne en RDC",
         0.4, 2.85, 12.5, 0.55,
         size=Pt(16), italic=True, color=RGBColor(0xAA, 0xCC, 0xFF),
         align=PP_ALIGN.CENTER)

# Boîte infos
add_rect(sl, 2.5, 3.7, 8.3, 2.0, fill_rgb=RGBColor(0x14, 0x38, 0x6E))
add_text(sl, "Présenté par :  KALUMBA ILUNGA Rooney",
         2.7, 3.82, 8, 0.42, size=Pt(15), bold=True, color=BLANC)
add_text(sl, "Dirigé par :  CT Ferdinand KAHENGA",
         2.7, 4.25, 8, 0.38, size=Pt(13), color=RGBColor(0xBB, 0xDD, 0xFF))
add_text(sl, "Filière : Génie Logiciel               Juin 2026",
         2.7, 4.65, 8, 0.38, size=Pt(13), color=RGBColor(0xBB, 0xDD, 0xFF))

add_text(sl, "Travail de Fin de Cycle – Ingénieur Technicien en Génie Logiciel",
         0.4, 6.7, 12.5, 0.4, size=Pt(10), italic=True,
         color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)

footer(sl, 1)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 – CONTEXTE & PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Contexte & Problématique",
           "La RDC face à la crise informationnelle")

# Colonne gauche – contexte
add_rect(sl, 0.25, 1.6, 6.0, 5.25, fill_rgb=BLANC)
add_text(sl, "Contexte congolais", 0.35, 1.7, 5.8, 0.45,
         size=Pt(14), bold=True, color=BLEU_FONCE)
add_bullet_box(sl, [
    "Explosion du mobile et de WhatsApp / Télégramme",
    "Médias en ligne : Radio Okapi, Actualite.cd, 7sur7.cd…",
    "Contenus non structurés, redondants, de qualité variable",
    "Événements sensibles : santé, sécurité, politique",
    "Contexte multilingue : français, anglais, swahili",
], 0.35, 2.15, 5.7, 3.5, size=Pt(13))

# Colonne droite – 2 phénomènes
add_rect(sl, 6.6, 1.6, 6.45, 2.4, fill_rgb=RGBColor(0xFF, 0xEE, 0xEE))
add_text(sl, "🔴  Désinformation", 6.75, 1.65, 6.2, 0.45,
         size=Pt(14), bold=True, color=RGBColor(0xCC, 0x00, 0x00))
add_bullet_box(sl, [
    "Contenus faux ou trompeurs viralisés",
    "Confiance accordée à l'expéditeur ≠ source réelle",
    "Messageries chiffrées : monitoring difficile",
], 6.75, 2.1, 6.1, 1.7, size=Pt(12), color=GRIS_TEXTE)

add_rect(sl, 6.6, 4.15, 6.45, 2.4, fill_rgb=RGBColor(0xFF, 0xF0, 0xCC))
add_text(sl, "🟠  Surinformation (infobésité)", 6.75, 4.2, 6.2, 0.45,
         size=Pt(14), bold=True, color=ORANGE)
add_bullet_box(sl, [
    "Volume > capacité de traitement cognitif",
    "Répétition = illusion de véracité",
    "Les outils de correction peuvent aggraver le bruit",
], 6.75, 4.65, 6.1, 1.7, size=Pt(12), color=GRIS_TEXTE)

# Question centrale
add_rect(sl, 0.25, 6.6, 12.83, 0.55, fill_rgb=BLEU_CLAIR)
add_text(sl, "❓ Comment concevoir un chatbot RAG transformant des articles d'actualité "
             "en réponses fiables, réduisant simultanément désinformation et surinformation sur WhatsApp/Télégramme ?",
         0.4, 6.63, 12.5, 0.48, size=Pt(11), bold=True, color=BLANC)
footer(sl, 2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 – OBJECTIFS & HYPOTHÈSES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Objectifs & Hypothèses", "Ce que le système vise à démontrer")

# Objectifs
add_rect(sl, 0.25, 1.6, 6.1, 5.1, fill_rgb=BLANC)
add_text(sl, "Objectifs du système", 0.4, 1.7, 5.8, 0.45,
         size=Pt(14), bold=True, color=BLEU_FONCE)
add_bullet_box(sl, [
    "Regrouper automatiquement les informations similaires",
    "Détecter et supprimer les redondances conversationnelles",
    "Produire une réponse synthétique, sourcée et structurée",
    "Extraire le texte des images via OCR (Tesseract)",
    "Déployer via WhatsApp (Whapi) et Télégramme Bot API",
    "Offrir un tableau de bord administrateur en temps réel",
], 0.4, 2.15, 5.8, 4.3, size=Pt(13))

# Hypothèse RAG
add_rect(sl, 6.6, 1.6, 6.45, 2.8, fill_rgb=RGBColor(0xE8, 0xF5, 0xE9))
add_text(sl, "Hypothèse centrale – Architecture RAG", 6.75, 1.68, 6.2, 0.5,
         size=Pt(13), bold=True, color=VERT)
add_bullet_box(sl, [
    "1. Analyse de la requête utilisateur",
    "2. Recherche sémantique (Top-K) dans le corpus",
    "3. Regroupement des contenus similaires",
    "4. Génération d'une réponse synthétique",
], 6.75, 2.2, 6.1, 1.9, size=Pt(12), color=GRIS_TEXTE, bullet="  ➤  ")

# Intérêt du sujet
add_rect(sl, 6.6, 4.55, 6.45, 2.15, fill_rgb=RGBColor(0xE3, 0xF2, 0xFD))
add_text(sl, "Intérêt du sujet", 6.75, 4.62, 6.2, 0.45,
         size=Pt(13), bold=True, color=BLEU_CLAIR)
add_bullet_box(sl, [
    "Scientifique : IR, NLP, RAG, LLM",
    "Technologique : crawler + embeddings + VectorDB + LLM local",
    "Sociétal : info fiable sur canaux déjà adoptés par les citoyens",
], 6.75, 5.08, 6.1, 1.45, size=Pt(12), color=GRIS_TEXTE)

footer(sl, 3)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 – ARCHITECTURE DU SYSTÈME (RAG)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Architecture du Système – Pipeline RAG",
           "De la réception du message au verdict structuré")

# Flèche pipeline – 5 étapes
steps = [
    ("📥", "Réception", "WhatsApp / Telegram\nWebhook (FastAPI)"),
    ("🔍", "Prétraitement", "OCR (Tesseract)\nTopic Gate\nDéduplication"),
    ("🧠", "RAG", "Embeddings\nChromaDB Top-K\nSeuil similarité"),
    ("💬", "Génération", "LLM Mistral 7B\n(Ollama local)\nPrompt structuré"),
    ("📤", "Restitution", "Verdict + sources\nAnti-redondance\nRéponse courte"),
]
box_w = 2.3
start_x = 0.25
for i, (icon, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + 0.12)
    add_rect(sl, x, 1.6, box_w, 3.6, fill_rgb=BLEU_FONCE)
    add_text(sl, icon, x, 1.75, box_w, 0.55,
             size=Pt(24), color=BLANC, align=PP_ALIGN.CENTER)
    add_text(sl, title, x, 2.3, box_w, 0.5,
             size=Pt(13), bold=True, color=RGBColor(0xAA, 0xDD, 0xFF),
             align=PP_ALIGN.CENTER)
    add_rect(sl, x + 0.05, 2.82, box_w - 0.1, 0.04, fill_rgb=BLEU_CLAIR)
    add_text(sl, desc, x + 0.1, 2.9, box_w - 0.2, 2.1,
             size=Pt(11), color=BLANC, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        ax = x + box_w + 0.02
        add_text(sl, "▶", ax, 2.9, 0.1, 0.5,
                 size=Pt(18), color=BLEU_CLAIR, align=PP_ALIGN.CENTER)

# Mémoire duale bas
add_rect(sl, 0.25, 5.45, 12.83, 1.35, fill_rgb=RGBColor(0x14, 0x38, 0x6E))
add_text(sl, "Mémoire duale — PostgreSQL (métadonnées structurées)  +  ChromaDB (vecteurs sémantiques)  +  Redis (état court terme)",
         0.5, 5.52, 12.3, 0.42, size=Pt(12), bold=True,
         color=RGBColor(0xAA, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
add_text(sl, "Synchronisation SQL → Chroma  ·  Relation 1:1 article_id  ·  Suppression doublons ON CONFLICT  ·  Couverture vectorielle 100 %",
         0.5, 5.94, 12.3, 0.38, size=Pt(11),
         color=RGBColor(0x88, 0xBB, 0xFF), align=PP_ALIGN.CENTER)

footer(sl, 4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 – FONCTIONNALITÉS CLÉS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Fonctionnalités Clés du Système",
           "Fact-checking + Anti-surinformation + OCR multicanal")

cards = [
    (BLEU_FONCE, "🔎  Fact-checking RAG",
     ["Retrieval sémantique Top-K",
      "Verdict : VRAI / FAUX / IMPRÉCIS / NON VÉRIFIABLE",
      "Sources citées (Radio Okapi, 7sur7, Actualite.cd…)",
      "Seuil de similarité configurable (RAG_MIN_SIMILARITY)"]),
    (VERT, "🔄  Anti-Surinformation",
     ["Embeddings conversationnels (chat_id)",
      "Cooldown par sujet : blocage des réponses répétées",
      "Story Card (fiche unique par cluster de messages)",
      "Mémoire Redis à TTL court (fenêtre glissante)"]),
    (ORANGE, "🖼️  OCR Multimodal",
     ["Détection automatique d'images contenant du texte",
      "Extraction via Tesseract OCR",
      "Injection dans le pipeline RAG",
      "Couvre captures, affiches, communiqués"]),
    (RGBColor(0x6A, 0x1B, 0x9A), "📊  Administration",
     ["Dashboard Web : corpus, sources, couverture",
      "17 595 articles indexés · 100 % vectorisés",
      "52 sources configurées (FR/EN/SW)",
      "API /admin/overview en temps réel"]),
]
card_w = 5.85
for i, (color, title, bullets) in enumerate(cards):
    row, col = divmod(i, 2)
    x = 0.25 + col * (card_w + 0.25)
    y = 1.6 + row * 2.65
    add_rect(sl, x, y, card_w, 2.5, fill_rgb=color)
    add_text(sl, title, x + 0.15, y + 0.1, card_w - 0.2, 0.5,
             size=Pt(14), bold=True, color=BLANC)
    add_bullet_box(sl, bullets, x + 0.15, y + 0.6, card_w - 0.3, 1.75,
                   size=Pt(12), color=BLANC, bullet_color=BLANC, bullet="· ")

footer(sl, 5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 – STACK TECHNIQUE & DÉPLOIEMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Stack Technique & Déploiement VPS",
           "Architecture locale, souveraine et opérationnelle")

# Tableau stack
cols = [
    ("Couche", BLEU_FONCE),
    ("Technologie", BLEU_FONCE),
    ("Rôle", BLEU_FONCE),
]
rows_data = [
    ("API & Webhooks",    "FastAPI (Python 3.12)",       "Endpoints WhatsApp, Télégramme, Admin"),
    ("LLM local",         "Ollama + Mistral 7B q4",      "Génération de réponses (0 coût API)"),
    ("Vecteurs",          "ChromaDB",                    "Recherche sémantique Top-K (cosine)"),
    ("Stockage SQL",      "PostgreSQL",                  "Articles, métadonnées, historique"),
    ("Cache & files",     "Redis",                       "Mémoire conversationnelle, cooldown"),
    ("OCR",               "Tesseract",                   "Extraction texte depuis images"),
    ("Process mgr",       "PM2",                         "Supervision des services applicatifs"),
    ("Reverse proxy",     "Nginx + HTTPS",               "Terminaison TLS, routage"),
    ("Crawler",           "Python (RSS / HTTP)",         "Ingestion continue 17 000+ articles"),
]
col_w = [2.8, 3.5, 5.9]
header_y = 1.62
row_h = 0.48
x_starts = [0.25, 3.1, 6.65]
# header
for j, (label, color) in enumerate(cols):
    add_rect(sl, x_starts[j], header_y, col_w[j] - 0.05, 0.45, fill_rgb=BLEU_FONCE)
    add_text(sl, label, x_starts[j]+0.1, header_y+0.05, col_w[j]-0.15, 0.35,
             size=Pt(12), bold=True, color=BLANC)
for i, (c0, c1, c2) in enumerate(rows_data):
    y = header_y + 0.45 + i * row_h
    bg = BLANC if i % 2 == 0 else RGBColor(0xEA, 0xF2, 0xFF)
    for j, (cell, xj, wj) in enumerate(zip([c0, c1, c2], x_starts, col_w)):
        add_rect(sl, xj, y, wj - 0.05, row_h - 0.03, fill_rgb=bg)
        add_text(sl, cell, xj+0.1, y+0.05, wj-0.2, row_h-0.1,
                 size=Pt(11), color=GRIS_TEXTE)

# Infos déploiement
add_rect(sl, 0.25, 6.0, 12.83, 0.95, fill_rgb=RGBColor(0x14, 0x38, 0x6E))
add_text(sl, "🖥️  VPS Linux unique  ·  Services sur 127.0.0.1  ·  "
             "Logs PM2 + FastAPI + Ollama  ·  URL prod : rooney-rdc.rooneykalumba.tech",
         0.5, 6.1, 12.3, 0.38, size=Pt(11), color=RGBColor(0xAA, 0xDD, 0xFF))
add_text(sl, "Incidents résolus : conflit port 8000, incohérence .env, transfert rsync, extensions DB manquantes",
         0.5, 6.5, 12.3, 0.35, size=Pt(10), italic=True, color=RGBColor(0x88, 0xBB, 0xFF))
footer(sl, 6)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 – RÉSULTATS & DÉMONSTRATIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Résultats Obtenus & Démonstrations",
           "Corpus, comportement en production, captures réelles")

# Métriques clés – 4 boîtes
metrics = [
    (BLEU_FONCE,  "18 595",   "articles indexés"),
    (VERT,        "100 %",    "couverture vectorielle"),
    (ORANGE,      "52",       "sources configurées"),
    (RGBColor(0x6A,0x1B,0x9A), "3 langues", "FR · EN · SW"),
]
for i, (color, val, label) in enumerate(metrics):
    x = 0.25 + i * 3.25
    add_rect(sl, x, 1.6, 3.0, 1.3, fill_rgb=color)
    add_text(sl, val, x, 1.65, 3.0, 0.75,
             size=Pt(30), bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, 2.38, 3.0, 0.45,
             size=Pt(12), color=BLANC, align=PP_ALIGN.CENTER)

# Scénarios de démonstration
add_rect(sl, 0.25, 3.1, 12.83, 0.42, fill_rgb=BLEU_FONCE)
add_text(sl, "Scénarios de démonstration validés", 0.4, 3.15, 12.5, 0.32,
         size=Pt(13), bold=True, color=BLANC)

scenarios = [
    ("WhatsApp  Groupe",  "Pipeline RAG complet → Verdict + Explication + Sources"),
    ("WhatsApp  Groupe",  "Message redondant → anti-surinformation, réponse courte"),
    ("WhatsApp  Privé",   "Question individuelle → réponse contextualisée"),
    ("WhatsApp  Privé",   "Image (OCR) → extraction texte → RAG → verdict"),
    ("Télégramme  Privé", "Même chaîne de vérification, interface Bot"),
    ("Admin  Web",        "Dashboard : corpus, sources, couverture 100 %"),
]
for i, (canal, comportement) in enumerate(scenarios):
    y = 3.55 + i * 0.51
    bg = BLANC if i % 2 == 0 else RGBColor(0xEA, 0xF2, 0xFF)
    add_rect(sl, 0.25, y, 3.5, 0.48, fill_rgb=bg)
    add_text(sl, canal, 0.35, y+0.06, 3.3, 0.36, size=Pt(11), bold=True, color=BLEU_FONCE)
    add_rect(sl, 3.8, y, 9.28, 0.48, fill_rgb=bg)
    add_text(sl, comportement, 3.9, y+0.06, 9.0, 0.36, size=Pt(11), color=GRIS_TEXTE)

footer(sl, 7)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 – DISCUSSION, CONCLUSION & PERSPECTIVES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_rgb=GRIS_CLAIR)
header_bar(sl, "Discussion, Conclusion & Perspectives",
           "Bilan critique et trajectoire future")

# Forces
add_rect(sl, 0.25, 1.6, 4.0, 4.5, fill_rgb=RGBColor(0xE8, 0xF5, 0xE9))
add_text(sl, "✅  Forces", 0.38, 1.67, 3.8, 0.45,
         size=Pt(13), bold=True, color=VERT)
add_bullet_box(sl, [
    "Architecture pragmatique et déployable",
    "Souveraineté des données (LLM local)",
    "Multi-canaux (WhatsApp + Télégramme)",
    "Déduplication conversationnelle efficace",
    "Bonne observabilité opérationnelle",
], 0.38, 2.13, 3.7, 3.7, size=Pt(12), color=GRIS_TEXTE, bullet_color=VERT)

# Limites
add_rect(sl, 4.5, 1.6, 4.0, 4.5, fill_rgb=RGBColor(0xFF, 0xEB, 0xEE))
add_text(sl, "⚠️  Limites", 4.63, 1.67, 3.8, 0.45,
         size=Pt(13), bold=True, color=RGBColor(0xCC, 0x00, 0x00))
add_bullet_box(sl, [
    "Dépendance CPU/RAM du VPS unique",
    "OCR sensible aux images dégradées",
    "Swahili encore sous-représenté",
    "Gestion complexe multi-transports simultanés",
    "Évaluation perceptuelle (pas de benchmark formel)",
], 4.63, 2.13, 3.7, 3.7, size=Pt(12), color=GRIS_TEXTE,
   bullet_color=RGBColor(0xCC, 0x00, 0x00))

# Perspectives
add_rect(sl, 8.75, 1.6, 4.33, 4.5, fill_rgb=RGBColor(0xE3, 0xF2, 0xFD))
add_text(sl, "🚀  Perspectives", 8.88, 1.67, 4.1, 0.45,
         size=Pt(13), bold=True, color=BLEU_CLAIR)
add_bullet_box(sl, [
    "Enrichissement du corpus swahili",
    "Résumé adaptatif selon contexte",
    "Métriques de « charge cognitive évitée »",
    "Séparation LLM / indexation sur nœuds dédiés",
    "Dashboard unifié (alertes, latence, erreurs)",
], 8.88, 2.13, 4.1, 3.7, size=Pt(12), color=GRIS_TEXTE, bullet_color=BLEU_CLAIR)

# Conclusion
add_rect(sl, 0.25, 6.3, 12.83, 0.7, fill_rgb=BLEU_FONCE)
add_text(sl, "\"Un assistant accessible sur les canaux que les citoyens utilisent déjà, "
             "capable de transformer un flot d'articles dispersés en réponses structurées, sourcées et moins redondantes.\"",
         0.4, 6.35, 12.5, 0.58, size=Pt(11), italic=True, color=BLANC)
footer(sl, 8)


# ── Sauvegarde ────────────────────────────────────────────────────────────────
out = "/home/rooney/Projects/PROJECT_PERSO/rdc-news-intelligence/RDC_News_Intelligence_Presentation.pptx"
prs.save(out)
print(f"✅  Fichier généré : {out}")
